"""Build the StefaNP design deck (pptx).

Reproduces: docs/StefaNP_design.pptx

Content is the architecture as actually built and measured this month, not a
wish-list: every number quoted appears in logs/ or docs/, and the failures are
on the slides beside the passes.

    /nfs/data/cxs1024/envs/demenv/bin/python make_design_deck.py
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / 'src'))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from hydropfn.paths import ROOT  # noqa: E402

INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTE = RGBColor(0x5A, 0x5F, 0x66)
ACC = RGBColor(0x0B, 0x5F, 0x8A)        # blue: structure
GOOD = RGBColor(0x1B, 0x7F, 0x4B)       # green: passed
BAD = RGBColor(0xB0, 0x3A, 0x2B)        # red: failed
WARN = RGBColor(0xB5, 0x7A, 0x0B)       # amber: unproven
BG = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def txt(slide, x, y, w, h, s, size=14, bold=False, color=INK,
        align=PP_ALIGN.LEFT, italic=False, space=4):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(s.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.italic = italic
        r.font.name = "Calibri"
    return tb


def box(slide, x, y, w, h, fill, line=None, radius=True):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(1.25)
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def labelled_box(slide, x, y, w, h, title, body, fill=WHITE, line=ACC,
                 tsize=13, bsize=10.5):
    box(slide, x, y, w, h, fill, line)
    txt(slide, x + 0.12, y + 0.07, w - 0.24, 0.3, title, tsize, True, ACC)
    if body:
        txt(slide, x + 0.12, y + 0.42, w - 0.24, h - 0.5, body, bsize, False,
            INK, space=2)


def arrow(slide, x1, y1, x2, y2, color=MUTE, width=1.6):
    from pptx.enum.shapes import MSO_CONNECTOR
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1),
                                   Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(width)
    return c


def slide(prs, title, sub=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bar = box(s, 0, 0, 13.333, 0.06, ACC, radius=False)  # noqa: F841
    txt(s, 0.5, 0.22, 12.4, 0.5, title, 26, True, INK)
    if sub:
        txt(s, 0.5, 0.78, 12.4, 0.4, sub, 13, False, MUTE)
    return s


def table(slide, x, y, w, rows, col_w, header=True, size=10.5, rh=0.3):
    """Simple text table using boxes -- pptx tables are fiddly to style."""
    yy = y
    for ri, row in enumerate(rows):
        xx = x
        for ci, cell in enumerate(row):
            cw = col_w[ci]
            if ri == 0 and header:
                box(slide, xx, yy, cw, rh, ACC, radius=False)
                txt(slide, xx + 0.06, yy + 0.03, cw - 0.12, rh, str(cell),
                    size, True, WHITE)
            else:
                fill = BG if ri % 2 else WHITE
                box(slide, xx, yy, cw, rh, fill, radius=False)
                col = INK
                t = str(cell)
                if t.startswith("PASS") or t.startswith("+"):
                    col = GOOD
                elif t.startswith("FAIL") or t.startswith("-"):
                    col = BAD
                elif t.startswith("?"):
                    col = WARN
                txt(slide, xx + 0.06, yy + 0.03, cw - 0.12, rh, t, size,
                    False, col)
            xx += cw
        yy += rh
    return yy


def build(out: Path) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    # ---------------------------------------------------------------- title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    box(s, 0, 0, 13.333, 7.5, ACC, radius=False)
    txt(s, 1.0, 2.3, 11.3, 1.0, "StefaNP", 54, True, WHITE)
    txt(s, 1.0, 3.3, 11.3, 0.9,
        "A multimodal masked autoencoder whose prediction is conditioned on a\n"
        "retrieved context set of sites — a Transformer Neural Process for hydrology",
        20, False, WHITE)
    txt(s, 1.0, 4.6, 11.3, 1.2,
        "Design, training of each component, how they connect,\n"
        "and how to judge whether each one earns its place",
        15, False, RGBColor(0xD6, 0xE6, 0xEF))
    txt(s, 1.0, 6.5, 11.3, 0.4,
        "Sibling to StefaLand (arXiv:2509.17942), extending it with cross-site "
        "in-context inference and DEM pixels   ·   2026-08-20",
        11, False, RGBColor(0xC0, 0xD8, 0xE6))

    # ------------------------------------------------------ why not "PFN"
    s = slide(prs, "Naming: this is a Neural Process, not a PFN",
              "Being precise about the mechanism keeps the claims honest")
    labelled_box(s, 0.5, 1.3, 4.0, 2.0, "PFN (TabPFN)",
                 "Prior-data fitted: pretrained on samples from a SYNTHETIC\n"
                 "prior to amortise Bayesian inference.\n\n"
                 "That is its defining property — not the in-context part.",
                 line=MUTE)
    labelled_box(s, 4.7, 1.3, 4.0, 2.0, "Neural Process / TNP",
                 "Context set -> predict at query.\n\n"
                 "TabPFN is a TNP with a synthetic prior.\n"
                 "The conditioning mechanism we want lives HERE.")
    labelled_box(s, 8.9, 1.3, 3.9, 2.0, "Time-series FMs",
                 "Moirai (any-variate, missing channels — closest),\n"
                 "Chronos, TimesFM, PatchTST / Ti-MAE.\n\n"
                 "Source of the temporal unit's recipe.", line=MUTE)
    box(s, 0.5, 3.6, 12.3, 1.15, BG, ACC)
    txt(s, 0.75, 3.75, 11.8, 0.9,
        "We pretrain on REAL masked data, so \"PFN\" would be inaccurate. It becomes "
        "defensible only if the\nsimulator-generated task prior is actually built — and "
        "even then it is a hybrid.", 14, False, INK)
    txt(s, 0.5, 5.0, 12.3, 1.4,
        "Working name  StefaNP  —  Stefa (family link to StefaLand, whose static unit we reuse)\n"
        "                                    +  NP (Neural Process: the honest mechanism)",
        16, True, ACC)

    # ------------------------------------------------- scope correction
    s = slide(prs, "What StefaLand already is, and what we actually add",
              "Corrected after reading the code on ICDS — the design is more incremental than first drawn")
    box(s, 0.5, 1.25, 6.0, 2.65, WHITE, MUTE)
    txt(s, 0.7, 1.36, 5.6, 0.3, "StefaLand (arXiv:2509.17942), as inspected",
        14, True, MUTE)
    txt(s, 0.7, 1.75, 5.6, 2.0,
        "Masked autoencoder that ALREADY FUSES static + time series\n"
        "  hidden 256 · 4 heads · 4 enc layers · d_ffd 512  (57 MB)\n"
        "  5 forcing series: P, RelHum, SWd, Tmax, Tmin\n"
        "  ~49 statics: GMTED terrain, HWSD + SoilGrid2 soils,\n"
        "     MSWX climate, land cover, catchsize, porosity, soil_depth\n"
        "  NO channel geometry of any kind\n"
        "  ships eval_probe.py (leave-region-out linear probe) and an\n"
        "     image pathway (zarr_image_loader, image_encoders/)",
        11, False, INK, space=1)
    box(s, 6.8, 1.25, 6.0, 2.65, WHITE, GOOD)
    txt(s, 7.0, 1.36, 5.6, 0.3, "What StefaNP genuinely adds", 14, True, GOOD)
    txt(s, 7.0, 1.75, 5.6, 2.0,
        "1. CROSS-SITE attention (the connector) — context from other\n"
        "    sites at inference. StefaLand is per-site only.\n"
        "2. DEM PIXELS as a modality (unit B) — StefaLand is\n"
        "    attribute-based, so it cannot work where no attribute\n"
        "    table exists, i.e. most of the globe.\n"
        "3. POINT MEASUREMENTS as tokens (unit D) — irregular\n"
        "    (var, value, covariate) visits; this is what delivers the\n"
        "    measured cross-variable gain.\n"
        "4. DECODERS for masked reconstruction in data space.",
        11, False, INK, space=1)
    box(s, 0.5, 4.1, 12.3, 1.0, BG, ACC)
    txt(s, 0.7, 4.22, 11.9, 0.85,
        "So the honest mapping is:  StefaLand  ≈  units A + C already fused.  An earlier slide split them as if "
        "the temporal\nunit were new — it is not. Good for feasibility (a pretrained backbone exists), "
        "sobering for novelty claims.", 12.5, False, INK)
    box(s, 0.5, 5.3, 12.3, 1.35, WHITE, WARN)
    txt(s, 0.7, 5.42, 11.9, 0.3, "Measured caveat before reuse", 13, True, WARN)
    txt(s, 0.7, 5.75, 11.9, 0.85,
        "Its embeddings reconstruct absolute elevation at R² 0.936 — but that is NOT learned memorisation: "
        "GMTED_elevation\nis one of its static INPUTS. De-locating is therefore a one-line ablation, not a "
        "research problem. They also added −0.004\nto channel geometry over 20 tabular attributes (expected: "
        "that task has a tabular shortcut).", 11.5, False, INK)

    # ------------------------------------------------- architecture (big)
    s = slide(prs, "Architecture: encoder / connector / decoder, three levels",
              "Every unit needs an ENCODER and a DECODER; the connector only ever sees site summaries")
    txt(s, 0.45, 1.03, 12.4, 0.3,
        "LEVEL 1 — per-site ENCODERS.  d = 256 everywhere (StefaLand's width, so its weights drop in unchanged).  "
        "B = batch, S = sites, M = visits.", 11, True, ACC)
    ux, uy, uw, uh = 0.45, 1.42, 2.9, 2.05
    units = [
        ("A · Static landscape", GOOD,
         "IN   (B,S,49) attributes\nENC  MLP 49->256->256\nOUT  (B,S,1,256)  1 token\n"
         "Pretrain: mask 15% of the 49\n  attrs, regress them back\nSTATUS reuse StefaLand"),
        ("B · Terrain (DEM)", GOOD,
         "IN   (B,S,1,128,128) @10 m\nENC  conv 128->16, pool\nOUT  (B,S,1..4,256)\n"
         "Pretrain: masked recon\n  (diffusion for the decoder)\nSTATUS built, sampler works"),
        ("C · Forcing-response", BAD,
         "IN   (B,S,T,V) T~3650 d, V~10\n  + the A and B tokens, in the\n  SAME sequence (that is how\n"
         "  attributes reach unit C)\nENC  16-d patches -> 228 tok\n     -> pool to 2-4 tokens\n"
         "OUT  (B,S,2..4,256)\nSTATUS NOT BUILT"),
        ("D · Measurement", GOOD,
         "IN   (B,S,M,3) var/val/cov\nENC  sum of 3 embeddings\nOUT  (B,S,M,256)  M~12 tok\n"
         "Pretrain: masked value\nSTATUS built, gates pass"),
    ]
    for i, (t, c, b) in enumerate(units):
        x = ux + i * (uw + 0.19)
        box(s, x, uy, uw, uh, WHITE, c)
        txt(s, x + 0.12, uy + 0.06, uw - 0.24, 0.3, t, 12.5, True, c)
        txt(s, x + 0.12, uy + 0.38, uw - 0.24, uh - 0.45, b, 9.5, False, INK,
            space=1)
        arrow(s, x + uw / 2, uy + uh, x + uw / 2, uy + uh + 0.3, c)
    box(s, 0.45, 3.8, 12.44, 0.52, BG, ACC)
    txt(s, 0.6, 3.88, 12.1, 0.4,
        "concat -> per-site bundle  (B, S, K, 256),  K = 4–10 tokens/site   ·   "
        "absent modality = absent token (no imputation)",
        12.5, True, ACC, PP_ALIGN.CENTER)
    arrow(s, 6.6, 4.32, 6.6, 4.6, ACC, 2.2)
    box(s, 0.45, 4.65, 12.44, 1.25, WHITE, ACC)
    txt(s, 0.6, 4.74, 12.1, 1.15,
        "LEVEL 2 — CONNECTOR, cross-site transformer   (8 layers, 8 heads, d=256, FFN 1024)\n"
        "IN    (B, 1 + S*K, 256)   flatten sites; prepend [TASK]; + geo-encoding; padding mask (B, 1+S*K)\n"
        "OUT   (B, 1 + S*K, 256)   -> take position 0 only\n"
        "permutation-invariant over sites (no positional index)  ·  query attends to context, never the reverse",
        11.5, True, INK, PP_ALIGN.CENTER, space=1)
    arrow(s, 6.6, 5.9, 6.6, 6.12, ACC, 2.2)
    box(s, 0.45, 6.16, 12.44, 0.72, WHITE, GOOD)
    txt(s, 0.6, 6.23, 12.1, 0.62,
        "LEVEL 3 — per-site DECODERS (one per modality).  IN [ context-aware summary (K,256) + unmasked patches ]"
        "  ->  OUT reconstructed data\n"
        "e.g. temporal decoder -> (B, N_masked, 16) days   ·   scalar head -> (B, 64) bin logits   ·   "
        "DEM decoder -> (B,1,128,128).   Loss is in DATA space.",
        11, True, GOOD, PP_ALIGN.CENTER, space=1)
    txt(s, 0.45, 6.95, 12.4, 0.4,
        "Why three levels: a 40-y daily record is ~900 patch tokens — too many for the connector, which must stay "
        "at 4–10 summary tokens/site. The summary is the CONDITIONING VECTOR that carries cross-site information "
        "down to the per-site decoder.", 10, False, MUTE)

    # -------------------------------------------------- the token sequence
    s = slide(prs, "What the final transformer actually reads",
              "One flat sequence; every token carries WHO it is, WHAT variable, and WHICH role")
    txt(s, 0.5, 1.15, 12.3, 0.3,
        "token  =  content embedding  +  variable-ID embedding  +  covariate embedding  "
        "+  role embedding  +  own/context flag", 13, True, ACC)
    rows = [
        ["#", "token", "content", "role", "own?"],
        ["0", "[TASK]", "var-ID(log_d) + cov(log_Q=3.2) + MASK-value", "query", "own"],
        ["1", "attr(query site)", "MLP(18 attributes)", "attribute", "own"],
        ["2", "dem(query site)", "terrain encoder(128x128 patch)", "attribute", "own"],
        ["3..k", "visit(query site)", "var-ID(log_W) + value + cov(log_Q)", "observed", "own"],
        ["k+1", "attr(context site 1)", "MLP(18 attributes)", "attribute", "context"],
        ["k+2", "ts(context site 1)", "pooled patches + FDC quantiles", "attribute", "context"],
        ["k+3..", "visit(context site 1)", "var-ID + value + cov", "observed", "context"],
        ["...", "... x N context sites", "(retrieved: attribute / geographic nearest)", "", "context"],
    ]
    table(s, 0.5, 1.6, 12.3, rows, [0.7, 2.5, 4.6, 1.7, 1.1], rh=0.34)
    box(s, 0.5, 4.85, 12.3, 1.1, BG, ACC)
    txt(s, 0.7, 4.97, 11.9, 0.9,
        "The [TASK] token is the query: it names the variable and the covariates but its VALUE is replaced by a\n"
        "learned mask embedding. The head reads out at position 0. Prediction = attention from [TASK] over everything else.",
        13, False, INK)
    txt(s, 0.5, 6.15, 12.3, 0.9,
        "Sizing:  200–500 context sites x ~5 tokens = 1–2.5k tokens.  Retrieval (which sites to include) is part of "
        "the method, not an afterthought.\nPhase-1 built version: 17 sites x (1 attr + 12 visits) = 221 tokens, "
        "1.2M parameters.", 12, False, MUTE)

    # ------------------------------------------------- how context is used
    s = slide(prs, "How context vectors are used — one objective, three regimes",
              "Masked-measurement modelling: hide a value, keep its variable-ID and covariates, predict it")
    regimes = [
        ("Own-site context", "the site's OTHER visits are visible", "at-a-station rating\n= amortised partial pooling",
         "A2 PASS: +0.081 / +0.167 / +0.081 R2\nover the per-site power law", GOOD),
        ("Cross-variable context", "only OTHER variables at this site", "infer depth where only\nwidth was measured",
         "A1 PASS: +0.108 R2 on velocity\n= 2.6x a hand-built RF feature", GOOD),
        ("No own context", "neighbouring sites only", "prediction in ungauged basins\n(the global mode)",
         "A3 FAIL on 2 of 3: worse than an\nattributes-only RF. The weak point.", BAD),
    ]
    for i, (t, ctx, use, res, c) in enumerate(regimes):
        x = 0.5 + i * 4.28
        box(s, x, 1.3, 4.0, 4.5, WHITE, c)
        txt(s, x + 0.15, 1.42, 3.7, 0.35, t, 15, True, c)
        txt(s, x + 0.15, 1.9, 3.7, 0.5, "context contains:", 10, True, MUTE)
        txt(s, x + 0.15, 2.15, 3.7, 0.6, ctx, 12, False, INK)
        txt(s, x + 0.15, 2.95, 3.7, 0.5, "use case:", 10, True, MUTE)
        txt(s, x + 0.15, 3.2, 3.7, 0.9, use, 12, False, INK)
        box(s, x + 0.15, 4.25, 3.7, 1.35, BG, None)
        txt(s, x + 0.3, 4.38, 3.4, 1.1, res, 11, False, c)
    txt(s, 0.5, 6.0, 12.3, 1.2,
        "Same weights, same forward pass — only which tokens are present changes. That is why ONE pretraining "
        "objective yields\nall three behaviours, and why the model degrades gracefully instead of failing when a "
        "modality is missing.", 13, False, INK)

    # --------------------------------------------- unit C: the time series
    s = slide(prs, "Unit C — how the time-series encoder is pretrained",
              "Four mask types; each one IS a downstream capability")
    rows = [
        ["mask", "what is hidden", "learn from", "capability unlocked"],
        ["random span", "contiguous chunk of an obs series", "forcings + surrounding obs", "gap filling"],
        ["causal tail", "the future", "past only (ONLY causal mask)", "forecasting"],
        ["whole-variable", "ALL of one series (e.g. soil moisture)", "forcings + the site's other series", "cross-variable inference"],
        ["whole-site", "every observation at the site", "forcings + attrs + NEIGHBOUR sites", "PUB (needs cross-site attn)"],
    ]
    table(s, 0.5, 1.3, 12.3, rows, [1.9, 3.7, 3.4, 3.3], rh=0.42)
    box(s, 0.5, 3.5, 6.0, 2.0, WHITE, ACC)
    txt(s, 0.65, 3.6, 5.7, 0.3, "Mechanics", 13, True, ACC)
    txt(s, 0.65, 3.95, 5.7, 1.5,
        "· daily series -> ~16-day patches -> tokens\n"
        "· channel-independent + variable-ID embedding\n"
        "· day-of-year features added\n"
        "· bar-distribution loss on masked patches\n"
        "· emits BOTH reconstructions and pooled summary tokens",
        11, False, INK, space=2)
    box(s, 6.8, 3.5, 6.0, 2.0, WHITE, WARN)
    txt(s, 6.95, 3.6, 5.7, 0.3, "The 'long record' problem", 13, True, WARN)
    txt(s, 6.95, 3.95, 5.7, 1.5,
        "40 y daily = 14,600 steps = ~900 patches per variable.\n"
        "Too many to hand to the connector.\n\n"
        "-> sample multi-year WINDOWS during training\n"
        "-> summarise the full record with FDC quantiles +\n"
        "    seasonal statistics beside pooled window tokens",
        11, False, INK, space=2)
    box(s, 0.5, 5.7, 12.3, 0.85, BG, GOOD)
    txt(s, 0.7, 5.82, 11.9, 0.7,
        "U3 gate for this unit: masked-span streamflow reconstruction must beat a REGIONAL LSTM on the same basins,\n"
        "under leave-region-out, before it is allowed into the connected model.", 13, True, GOOD)

    # ------------------------------- unit C, the three design questions
    s = slide(prs, "Unit C — three design questions, answered",
              "Raised 2026-08-20; these change the pretraining recipe")
    qa = [
        ("Q  Doesn't this need tokens from the OTHER units to be done properly?",
         "YES — and that reshapes the plan. Streamflow from forcings alone is under-determined; the\n"
         "response depends on catchment properties that live in units A and B. So unit C is pretrained in\n"
         "TWO stages: (1) ALONE, forcings->obs, to prove the temporal encoder works and clear its U3 gate;\n"
         "(2) JOINTLY, cross-attending to the A/B tokens of the SAME site, which is where basin identity\n"
         "enters. Stage 1 is a scaffold, not the deliverable — and stage 2's gain over stage 1 is exactly\n"
         "the U4 evidence that the other units carry something.", GOOD),
        ("Q  Predict observations from forcings, or also infer forcings from observations?",
         "BOTH — masking must be BIDIRECTIONAL over channels. Forcing->obs alone teaches a one-way map;\n"
         "masking precipitation and recovering it from streamflow is a real inverse problem (and a real\n"
         "product: precipitation correction in poorly-gauged basins). Implementation is free: the mask\n"
         "sampler simply does not distinguish 'forcing' from 'observation' channels. Every channel is a\n"
         "token stream with a variable-ID; any of them can be hidden.", GOOD),
        ("Q  When pretraining the whole thing, do we mask attributes and predict others?",
         "YES, and ACROSS modalities, not just within them. The connected objective masks any subset of:\n"
         "static attributes  ·  DEM patches  ·  observation series  ·  forcing series  ·  point measurements\n"
         "— and reconstructs them from whatever remains, including other SITES. That single cross-modal\n"
         "mask is what makes one model serve every task, and it is why missing-modality dropout during\n"
         "training is mandatory rather than optional.", GOOD),
    ]
    y = 1.2
    for q, a, c in qa:
        box(s, 0.5, y, 12.3, 1.78, WHITE, c)
        txt(s, 0.7, y + 0.08, 11.9, 0.3, q, 13, True, c)
        txt(s, 0.7, y + 0.42, 11.9, 1.3, a, 10.5, False, INK, space=1)
        y += 1.92
    txt(s, 0.5, y + 0.02, 12.3, 0.4,
        "Consequence: 'pretrain each unit separately, then connect' is right for DEBUGGING, but the final "
        "objective must be joint and cross-modal.", 11.5, True, ACC)

    # ----------------------------------------- variable identity (the Q)
    s = slide(prs, "Are variables distinguished by their index in the vector?",
              "No — and that choice is what makes missing and new variables work")
    box(s, 0.5, 1.3, 6.0, 3.2, WHITE, BAD)
    txt(s, 0.7, 1.42, 5.6, 0.35, "Fixed-index / channel approach", 15, True, BAD)
    txt(s, 0.7, 1.9, 5.6, 2.4,
        "x = [ Q, SM, SWE, T, ... ]\n"
        "variable identity = its SLOT in the vector\n\n"
        "· schema fixed at training time\n"
        "· a missing variable must be imputed or zero-filled,\n"
        "   and the model cannot tell 'zero' from 'absent'\n"
        "· adding a new variable = retrain everything\n"
        "· cannot handle a site that measures a different subset",
        12, False, INK, space=3)
    box(s, 6.8, 1.3, 6.0, 3.2, WHITE, GOOD)
    txt(s, 7.0, 1.42, 5.6, 0.35, "Variable-ID embedding (ours, and Moirai's)", 15, True, GOOD)
    txt(s, 7.0, 1.9, 5.6, 2.4,
        "token = var_emb[id] + value_emb(v) + cov_emb(Q,t)\n"
        "variable identity = a LEARNED VECTOR it carries\n\n"
        "· any subset of variables, any order\n"
        "· absent variable = absent token (unambiguous)\n"
        "· new variable = add an embedding row, fine-tune\n"
        "· one head serves all variables; the bar-distribution\n"
        "   borders are per-variable",
        12, False, INK, space=3)
    box(s, 0.5, 4.75, 12.3, 1.05, BG, ACC)
    txt(s, 0.7, 4.87, 11.9, 0.85,
        "Consequence for the connector: the sequence is a SET, not a fixed record. Permutation invariance over both "
        "sites and\nvariables is what lets one model serve every task — and it is why 'missing data' needs no "
        "imputation anywhere.", 13, False, INK)
    txt(s, 0.5, 6.0, 12.3, 0.9,
        "Built and verified in Phase 1: nn.Embedding(N_VARS, d) added to every measurement token; the [TASK] token "
        "carries the\nqueried variable's ID with a learned MASK vector in place of the value.  (lib/hydropfn.py)",
        12, False, MUTE)

    # ------------------------------------------- judging each component
    s = slide(prs, "How to judge whether a component is useful",
              "Five steps; U0–U3 run WITHOUT the connected model, so a dead unit is found in days")
    rows = [
        ["step", "question", "gate it must clear"],
        ["U0 adapter", "can raw data become (site, time, var, value, cov) records?", "round-trips; missingness explicit, never imputed"],
        ["U1 encoder", "can it become a small set of tokens?", "fixed token count; absent modality handled"],
        ["U2 objective", "can it be self-supervised alone?", "loss falls on held-out REGIONS, not just rows"],
        ["U3 standalone", "does the representation beat a simple baseline?", "named baseline, leave-region-out, noise ceiling measured FIRST"],
        ["U4 connection", "does it add anything the other units lack?", "must move a downstream metric; DROP it and see"],
    ]
    table(s, 0.5, 1.3, 12.3, rows, [1.9, 5.6, 4.8], rh=0.46)
    box(s, 0.5, 4.15, 12.3, 1.5, WHITE, BAD)
    txt(s, 0.7, 4.27, 11.9, 0.3, "Why U4 exists — measured three separate times", 14, True, BAD)
    txt(s, 0.7, 4.62, 11.9, 0.95,
        "StefaLand embeddings −0.004    ·    upstream network structure ±0.001    ·    cross-section CNN ~0\n\n"
        "A unit can encode real physics and still add NOTHING if simpler inputs already carry it. "
        "Passing U3 is not permission to ship.", 13, False, INK)
    box(s, 0.5, 5.85, 12.3, 1.15, BG, ACC)
    txt(s, 0.7, 5.97, 11.9, 0.95,
        "Standing rules: split by SITE never COMID  ·  identical row/site population across compared models "
        "(a forgotten\ncontrol turned a real +0.024 into an apparent +0.090)  ·  report SEEDS not single runs  ·  "
        "exercise a metric on a known input before believing a table built from it.", 12, False, INK)

    # ------------------------------------------------- evidence so far
    s = slide(prs, "Evidence so far — what is proven, what failed",
              "Every number from a leave-region-out run this month; failures kept on the slide")
    rows = [
        ["gate", "what it tests", "result", "verdict"],
        ["T2 cross-variable (RF)", "hand-built width anomaly -> velocity", "+0.041 R2", "PASS"],
        ["T1 at-a-station (TabPFN)", "in-context vs train-once, same info", "+0.087 / +0.056 median-site", "PASS"],
        ["StefaNP A1", "cross-variable, no hand-built feature", "+0.108 R2 = 2.6x the feature", "PASS"],
        ["StefaNP A2", "beat per-site power law", "+0.081 / +0.167 / +0.081", "PASS"],
        ["StefaNP A3", "zero-context vs attributes-only RF", "-0.075 / +0.013 / -0.042", "FAIL 2 of 3"],
        ["Sampler v0 (uncond.)", "RePaint forcing on unconditional DDPM", "psd 2.90 vs band 0.7-1.3", "FAIL"],
        ["Sampler v1 (conditioned)", "texture + diversity + fidelity", "psd 0.66, vario 0.80/0.93, elev 1.05", "PASS"],
        ["Terrain unit U4", "does DEM add to the connected model?", "not yet run", "?  unproven"],
        ["Unit C (time series)", "beat a regional LSTM", "not built", "?  unproven"],
    ]
    table(s, 0.5, 1.25, 12.3, rows, [3.0, 4.3, 3.5, 1.5], rh=0.4)
    box(s, 0.5, 5.35, 12.3, 1.05, BG, WARN)
    txt(s, 0.7, 5.47, 11.9, 0.85,
        "Caveat that applies to ALL of the above: single seed, single holdout region (HUC2 03). This session's own "
        "history says\nsingle runs mislead by more than the effects being claimed — multi-seed, multi-region before "
        "any of this is quoted.", 12.5, False, INK)
    txt(s, 0.5, 6.55, 12.3, 0.5,
        "Also open: a join defect leaves site_no NaN on 60% of rows, so Phase 1 trains on 25,656 of 64,797 visits.",
        12, False, BAD)

    # --------------------------------- the mask IS the query
    s = slide(prs, "One model, every prediction mode: the MASK IS THE QUERY",
              "There is no mode switch at inference — you choose which positions to hide, and that is the question")
    rows = [
        ["your data situation", "what you mask at inference", "= which training mask", "what you get"],
        ["forcings only, no observations ever", "all observation positions", "whole-site", "PUB / ungauged prediction"],
        ["some observations exist", "only the ones you want predicted;\nthe rest stay VISIBLE", "whole-site (partial)",
         "prediction conditioned on\nthe observations you have"],
        ["record has gaps", "the gap positions", "random span", "gap filling"],
        ["want the future", "everything after t", "causal tail", "forecasting"],
        ["one variable never measured", "that whole series", "whole-variable", "cross-variable inference"],
    ]
    table(s, 0.5, 1.25, 12.3, rows, [3.5, 3.4, 2.5, 2.9], rh=0.62, size=10.5)
    box(s, 0.5, 4.5, 12.3, 1.25, WHITE, GOOD)
    txt(s, 0.7, 4.6, 11.9, 0.3,
        "The key point: existing observations do NOT go into a separate 'context' slot",
        13.5, True, GOOD)
    txt(s, 0.7, 4.95, 11.9, 0.75,
        "Being UNMASKED *is* being context. A site with 3 years of streamflow and a gap in year 2 has: 3 years "
        "unmasked (context)\nand the gap masked (query) — in the same sequence, in one forward pass. The four "
        "training masks exist so that every\ncombination of visible and hidden the real world produces has been "
        "seen during training.", 11.5, False, INK)
    box(s, 0.5, 5.9, 12.3, 1.05, BG, WARN)
    txt(s, 0.7, 6.0, 11.9, 0.9,
        "Implementation trap: THREE kinds of 'absent' must never be conflated —  (1) padding, a short record: "
        "masked, never scored;\n(2) genuinely missing, a gauge outage: masked, never scored, but CAN be predicted; "
        "(3) deliberately masked for training:\npredicted AND scored. Filling any of them with zeros makes 'absent' "
        "indistinguishable from a real zero.", 11.5, False, INK)

    # --------------------------------- how to train it
    s = slide(prs, "How to train the system",
              "Four stages; each one is runnable and checkable before the next begins")
    stages = [
        ("0", "Unit-wise pretraining  —  DEBUGGING ONLY", GOOD,
         "Each encoder alone on its own masked objective. Purpose is to prove the encoder works and clear its\n"
         "U3 gate, NOT to produce the final weights. Unit B done (sampler); unit D done (Phase 1); A is StefaLand."),
        ("1", "Per-site joint  —  all modalities in ONE sequence", ACC,
         "sequence = [ attr_tok, dem_tok, meas_tok x M, ts_patch x N ].  Self-attention lets every temporal patch\n"
         "see the catchment descriptors — this is how attributes reach unit C, and it is what StefaLand already does.\n"
         "Objective: mask any subset of ANY modality, reconstruct in data space. Cross-modal from the start."),
        ("2", "Cross-site  —  add the connector", ACC,
         "Freeze or LoRA the per-site encoders; train the connector on task-sampled batches.\n"
         "WHOLE-SITE masking is what forces it to use neighbours: with every observation at the query site hidden,\n"
         "the only remaining information is other sites. Randomise the CONTEXT SIZE here (measured problem: with a\n"
         "fixed n_ctx the model calibrates off one neighbour and never learns to aggregate)."),
        ("3", "Joint fine-tune  —  everything, low LR", WARN,
         "Unfreeze with residual adapters (StefaLand ships these). Missing-modality dropout throughout so any\n"
         "subset of inputs works at inference. This is the only stage needing the 8xA100 machine."),
    ]
    y = 1.2
    for n, t, c, b in stages:
        h = 1.32 if n in ("2", "1") else 1.0
        box(s, 0.5, y, 0.5, h, c)
        txt(s, 0.5, y + h / 2 - 0.22, 0.5, 0.4, n, 19, True, WHITE,
            PP_ALIGN.CENTER)
        box(s, 1.08, y, 11.75, h, WHITE, c)
        txt(s, 1.25, y + 0.06, 11.4, 0.3, t, 12.5, True, c)
        txt(s, 1.25, y + 0.36, 11.4, h - 0.42, b, 10, False, INK, space=1)
        y += h + 0.13
    box(s, 0.5, y + 0.02, 12.3, 1.02, BG, ACC)
    txt(s, 0.7, y + 0.12, 11.9, 0.9,
        "Batch = a TASK, not a row: sample S sites, sample a mask pattern, mask, reconstruct. Loss = sum of "
        "per-modality\nreconstruction in data space, weighted so no modality dominates by sheer token count. "
        "Mask ratio on a schedule\n(easy -> hard). Splits by SITE and leave-region-out throughout — the only "
        "honest test of what this is for.", 11.5, False, INK)

    # ------------------------------- channel geometry into StefaLand
    s = slide(prs, "StefaLand has no channel geometry — three ways to add it",
              "Ordered by cost; they are complementary, not alternatives")
    ways = [
        ("1 · As an auxiliary TARGET  (shapes the encoder)", GOOD,
         "Add a head predicting log_d / log_v from the site's static+DEM representation.\n"
         "The head is DISCARDED at inference, so no geometry is needed at test time —\n"
         "it exists only to force the encoder to represent hydrologic context.\n"
         "MEASURED: raw log_d/log_v are learnable at R2 0.68–0.72, so the signal is dense.\n"
         "WARNING (measured): the RESIDUAL is NOT learnable from terrain — all six cells\n"
         "negative under two fold structures. Use raw targets, never residuals."),
        ("2 · As context TOKENS  (the StefaNP route, already built)", GOOD,
         "Each (var, value, Q) visit becomes a token in the site bundle.\n"
         "Needs no schema change to StefaLand — it is a new unit beside it.\n"
         "MEASURED: this is what delivers A1 (+0.108) and A2."),
        ("3 · As a fine-tuning corpus  (residual adapters)", WARN,
         "StefaLand ships residual fine-tuning adapters — the intended mechanism.\n"
         "Cheapest to try, but inherits StefaLand's location-awareness."),
    ]
    y = 1.25
    for t, c, b in ways:
        h = 1.75 if "MEASURED: raw" in b else 1.15
        box(s, 0.5, y, 12.3, h, WHITE, c)
        txt(s, 0.7, y + 0.08, 11.9, 0.3, t, 14, True, c)
        txt(s, 0.7, y + 0.42, 11.9, h - 0.5, b, 11.5, False, INK, space=1)
        y += h + 0.16
    box(s, 0.5, y, 12.3, 0.95, BG, BAD)
    txt(s, 0.7, y + 0.08, 11.9, 0.8,
        "Caveat carried from our measurements: StefaLand embeddings reconstruct absolute elevation at R2 0.936 — a "
        "LOCATION FINGERPRINT.\nThey added −0.004 to channel geometry. Evaluate under leave-region-out and consider a "
        "de-located variant before trusting transfer.", 11.5, False, INK)

    # ------------------------------- applying in a new region
    s = slide(prs, "Applying the model in a NEW region with some context data",
              "No retraining, no fine-tuning — the adaptation is the forward pass")
    steps = [
        ("1", "Assemble the query site", "attributes + DEM patch + whatever forcing series exist.\nNo observations required."),
        ("2", "RETRIEVE context sites", "top-k by attribute / geographic similarity from anywhere in the world.\nThey need not be in the new region — that is the point of a learned prior."),
        ("3", "Add any local observations", "a few gauges in the region, or two field visits at the query site itself.\nThey enter as ordinary context tokens; zero is allowed."),
        ("4", "Emit the [TASK] token", "name the variable and covariates (e.g. streamflow on date t; or log_d at Q).\nValue replaced by the learned mask embedding."),
        ("5", "Read the distribution", "one forward pass -> full predictive distribution, not a point estimate."),
    ]
    y = 1.2
    for n, t, b in steps:
        box(s, 0.5, y, 0.55, 0.92, ACC)
        txt(s, 0.5, y + 0.22, 0.55, 0.4, n, 20, True, WHITE, PP_ALIGN.CENTER)
        box(s, 1.15, y, 11.65, 0.92, WHITE, ACC)
        txt(s, 1.32, y + 0.06, 11.3, 0.3, t, 13.5, True, INK)
        txt(s, 1.32, y + 0.36, 11.3, 0.55, b, 11, False, MUTE, space=1)
        y += 1.02
    box(s, 0.5, y + 0.05, 12.3, 1.0, BG, GOOD)
    txt(s, 0.7, y + 0.16, 11.9, 0.85,
        "Time-series prediction specifically: the query is a [TASK] token per (variable, timestep) — or a masked "
        "PATCH for a horizon.\nContext = neighbouring gauged basins' forcing+discharge histories. This is the "
        "'whole-site mask' the temporal unit was pretrained on.", 12.5, True, GOOD)

    # ------------------------------------------------------- roadmap
    s = slide(prs, "What is next, and when the big machine gets used",
              "Small-scale gates first; 8xA100 only after value is confirmed")
    rows = [
        ["#", "next step", "why now", "cost"],
        ["1", "Fix the site_no join defect", "recovers 60% of training visits; helps every later number", "hours"],
        ["2", "Multi-seed / multi-region Phase 1", "current results are single-seed, single-holdout", "1 GPU-day"],
        ["3", "Build unit C (forcing-response)", "the only unbuilt unit; makes this a hydrology model", "2–5 GPU-days"],
        ["4", "Terrain U4 (E0/E1/E2)", "does DEM add anything the core lacks?", "1–2 GPU-days"],
        ["5", "Connection phase at scale", "all units, missing-modality dropout, adapters", "1–2 GPU-weeks"],
    ]
    table(s, 0.5, 1.3, 12.3, rows, [0.6, 4.2, 5.6, 1.9], rh=0.44)
    box(s, 0.5, 4.05, 12.3, 1.15, WHITE, BAD)
    txt(s, 0.7, 4.17, 11.9, 0.3, "The weakness to attack first", 14, True, BAD)
    txt(s, 0.7, 4.5, 11.9, 0.7,
        "A3 — the ZERO-CONTEXT floor — is the mode global/ungauged deployment actually runs in, and it is the one "
        "gate that failed.\nSteps 1 and 2 are aimed squarely at it before any new unit is added.", 12.5, False, INK)
    box(s, 0.5, 5.4, 12.3, 1.0, BG, ACC)
    txt(s, 0.7, 5.52, 11.9, 0.85,
        "Compute gating (agreed): the 8xA100 machine is engaged only after the single-GPU gates confirm value. "
        "Total pretraining\ncost is academic-scale (~2–4 GPU-weeks single-GPU equivalent); the real cost is data "
        "engineering, measured in person-weeks.", 12.5, False, INK)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default=str(ROOT / "docs" / "StefaNP_design.pptx"))
    build(Path(ap.parse_args().out))
